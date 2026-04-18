"""
Document Loader — parses files from the docs/ folder with rich metadata extraction.
Extracts page numbers, section headings, version, and effective date from filenames and content.
"""

import csv
import json
import os
import re
from pathlib import Path
from typing import List, Dict, Any, Tuple

DOCS_DIR = Path(__file__).parent / "docs"
CHUNK_SIZE = 600
CHUNK_OVERLAP = 100


# ── Metadata extraction from filename ────────────────────────────────────────

def _extract_file_metadata(path: Path) -> Dict[str, str]:
    """Extract version and effective date from filename patterns."""
    name = path.stem
    version = ""
    effective_date = ""

    # Version: v1, v2.1, version_3, _v3 etc.
    v_match = re.search(r'[_\-\s]v(\d+[\d.]*)', name, re.IGNORECASE)
    if not v_match:
        v_match = re.search(r'version[_\-\s]*(\d+[\d.]*)', name, re.IGNORECASE)
    if v_match:
        version = f"v{v_match.group(1)}"

    # Date: 2024-07-01, 2024_07, 20240701 etc.
    d_match = re.search(r'(\d{4}[-_]\d{2}[-_]\d{2})', name)
    if not d_match:
        d_match = re.search(r'(\d{4}[-_]\d{2})', name)
    if d_match:
        effective_date = d_match.group(1).replace("_", "-")

    return {"version": version, "effective_date": effective_date}


def _detect_section(text: str) -> str:
    """Try to detect a section heading at the start of a chunk."""
    lines = text.strip().split("\n")
    for line in lines[:3]:
        line = line.strip()
        # Looks like a heading: short, possibly numbered, not ending with period
        if line and len(line) < 120 and not line.endswith("."):
            if re.match(r'^(\d+[\.\d]*\s|[A-Z][A-Z\s]{3,}$|#{1,4}\s)', line):
                return line[:80]
    return ""


# ── Per-format extractors (with page tracking) ────────────────────────────────

def _extract_pdf(path: Path) -> List[Tuple[str, int]]:
    """Returns list of (text, page_number) tuples."""
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append((text, i))
    return pages


def _extract_docx(path: Path) -> List[Tuple[str, int]]:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return [("\n".join(parts), None)]


def _extract_xlsx(path: Path) -> List[Tuple[str, int]]:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), data_only=True)
    pages = []
    for i, sheet in enumerate(wb.worksheets, 1):
        rows = []
        rows.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(v) for v in row if v is not None)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            pages.append(("\n".join(rows), i))
    return pages


def _extract_pptx(path: Path) -> List[Tuple[str, int]]:
    from pptx import Presentation
    prs = Presentation(str(path))
    pages = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if len(parts) > 1:
            pages.append(("\n".join(parts), i))
    return pages


def _extract_csv(path: Path) -> List[Tuple[str, int]]:
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        text = "\n".join(" | ".join(row) for row in reader if any(c.strip() for c in row))
    return [(text, None)]


def _extract_txt(path: Path) -> List[Tuple[str, int]]:
    return [(path.read_text(encoding="utf-8", errors="ignore"), None)]


def _extract_json(path: Path) -> List[Tuple[str, int]]:
    def flatten(obj, prefix=""):
        lines = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                label = f"{prefix}{k}" if not prefix else f"{prefix} — {k}"
                lines.extend(flatten(v, label))
        elif isinstance(obj, list):
            for item in obj:
                lines.extend(flatten(item, prefix))
        else:
            val = str(obj).strip()
            if val:
                lines.append(f"{prefix}: {val}" if prefix else val)
        return lines
    data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    return [("\n".join(flatten(data)), None)]


EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".doc":  _extract_docx,
    ".xlsx": _extract_xlsx,
    ".xls":  _extract_xlsx,
    ".pptx": _extract_pptx,
    ".ppt":  _extract_pptx,
    ".csv":  _extract_csv,
    ".txt":  _extract_txt,
    ".md":   _extract_txt,
    ".json": _extract_json,
}


# ── Chunking ──────────────────────────────────────────────────────────────────

def _chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        if end < len(text):
            for boundary in ["\n\n", "\n", ". ", "! ", "? "]:
                idx = text.rfind(boundary, start, end)
                if idx > start:
                    end = idx + len(boundary)
                    break
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap if end - overlap > start else end
    return chunks


# ── Public API ────────────────────────────────────────────────────────────────

def load_documents_from_folder(folder: Path = DOCS_DIR) -> List[Dict[str, Any]]:
    folder.mkdir(parents=True, exist_ok=True)
    results = []

    for path in sorted(folder.iterdir()):
        if not path.is_file():
            continue
        ext = path.suffix.lower()
        extractor = EXTRACTORS.get(ext)
        if extractor is None:
            continue

        try:
            pages = extractor(path)
        except Exception as e:
            print(f"[loader] Could not parse {path.name}: {e}")
            continue

        file_meta = _extract_file_metadata(path)

        chunk_idx = 0
        for page_text, page_num in pages:
            chunks = _chunk_text(page_text)
            total = len(chunks)
            for i, chunk in enumerate(chunks):
                section = _detect_section(chunk)
                page_label = f"p.{page_num}" if page_num else ""
                results.append({
                    "id":               f"{path.stem}-chunk{chunk_idx}",
                    "source":           path.name,
                    "page":             page_num,
                    "page_label":       page_label,
                    "section":          section or f"Section {chunk_idx + 1}",
                    "title":            path.stem.replace("_", " ").replace("-", " "),
                    "version":          file_meta["version"],
                    "effective_date":   file_meta["effective_date"],
                    "content":          chunk,
                    "similarity_score": 0.0,
                })
                chunk_idx += 1

    return results


def save_uploaded_file(filename: str, data: bytes) -> Path:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    dest = DOCS_DIR / filename
    dest.write_bytes(data)
    return dest
